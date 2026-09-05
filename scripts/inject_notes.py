#!/usr/bin/env python3
"""Inject notes-only JSON into a copy of an existing PowerPoint deck."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def set_notes(slide, text: str) -> None:
    frame = slide.notes_slide.notes_text_frame
    lines = text.splitlines()
    frame.text = lines[0] if lines else ""
    for line in lines[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = line


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--notes", required=True, type=Path)
    parser.add_argument("--mode", choices=("replace", "append", "skip-if-present"), default="replace")
    args = parser.parse_args()

    if args.input.resolve() == args.output.resolve():
        raise SystemExit("Refusing to overwrite the source deck. Choose a different output path.")

    entries = json.loads(args.notes.read_text(encoding="utf-8"))
    notes = {int(item["slide"]): str(item["notes"]).strip() for item in entries}
    presentation = Presentation(str(args.input))
    if sorted(notes) != list(range(1, len(presentation.slides) + 1)):
        raise SystemExit("Notes JSON must contain exactly one entry for every slide.")

    for index, slide in enumerate(presentation.slides, start=1):
        new_text = notes[index]
        existing = (slide.notes_slide.notes_text_frame.text or "").strip()
        if args.mode == "skip-if-present" and existing:
            continue
        if args.mode == "append" and existing:
            new_text = f"{existing}\n\n{new_text}"
        set_notes(slide, new_text)

    args.output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(args.output))
    print(json.dumps({"slides": len(presentation.slides), "output": str(args.output)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
