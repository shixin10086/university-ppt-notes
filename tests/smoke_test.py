#!/usr/bin/env python3
"""End-to-end smoke test for JSON, DOCX, and PPTX export scripts."""

from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path

from docx import Document
from pptx import Presentation


def run(*args: str) -> None:
    subprocess.run([sys.executable, *args], check=True, capture_output=True, text=True, encoding="utf-8")


def normalize(text: str) -> str:
    return text.replace("\r\n", "\n").replace("\r", "\n").strip()


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    sample = root / "examples" / "sample-notes.md"
    with tempfile.TemporaryDirectory(prefix="university-ppt-notes-") as temp:
        temp_dir = Path(temp)
        source_pptx = temp_dir / "source.pptx"
        output_pptx = temp_dir / "with-notes.pptx"
        notes_json = temp_dir / "notes.json"
        output_docx = temp_dir / "notes.docx"
        state_json = temp_dir / ".university-ppt-notes" / "state.json"

        presentation = Presentation()
        while len(presentation.slides) < 5:
            presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(source_pptx)

        run(str(root / "scripts" / "build_notes_json.py"), "--input", str(sample), "--output", str(notes_json))
        run(str(root / "scripts" / "export_docx.py"), "--input", str(sample), "--output", str(output_docx))
        run(
            str(root / "scripts" / "inject_notes.py"),
            "--input", str(source_pptx),
            "--output", str(output_pptx),
            "--notes", str(notes_json),
        )
        run(
            str(root / "scripts" / "workflow_state.py"),
            "init",
            "--state", str(state_json),
            "--ppt", str(source_pptx),
            "--slide-count", "5",
        )
        run(
            str(root / "scripts" / "workflow_state.py"),
            "update",
            "--state", str(state_json),
            "--stage", "batch_written",
            "--confirmed-through", "5",
        )

        expected = {int(item["slide"]): normalize(item["notes"]) for item in json.loads(notes_json.read_text(encoding="utf-8"))}
        result = Presentation(output_pptx)
        for index, slide in enumerate(result.slides, start=1):
            actual = normalize(slide.notes_slide.notes_text_frame.text)
            if actual != expected[index]:
                raise AssertionError(f"Notes mismatch on slide {index}")
            if actual.startswith(f"P{index:02d}"):
                raise AssertionError(f"Page heading leaked into notes on slide {index}")

        document = Document(output_docx)
        text = "\n".join(p.text for p in document.paragraphs)
        for index in range(1, 6):
            if f"P{index:02d}" not in text:
                raise AssertionError(f"DOCX missing slide heading P{index:02d}")

        state = json.loads(state_json.read_text(encoding="utf-8"))
        if state["confirmed_through"] != 5 or state["stage"] != "batch_written":
            raise AssertionError("Workflow state did not preserve progress metadata")
        forbidden_state_keys = {"requirements", "instructions", "slide_text", "notes_content"}
        if forbidden_state_keys.intersection(state):
            raise AssertionError("Workflow state contains forbidden content fields")

        resumed_state = temp_dir / "resumed-state.json"
        run(
            str(root / "scripts" / "workflow_state.py"),
            "init",
            "--state", str(resumed_state),
            "--ppt", str(source_pptx),
            "--slide-count", "10",
            "--target-start", "6",
            "--confirmed-through", "5",
        )
        resumed = json.loads(resumed_state.read_text(encoding="utf-8"))
        if resumed["current_batch"] != {"start": 6, "end": 10} or resumed["confirmed_through"] != 5:
            raise AssertionError("Workflow state could not resume from a later target range")

    print(json.dumps({"slides": 5, "notes_match": True, "docx_opened": True, "state_ok": True, "status": "passed"}))


if __name__ == "__main__":
    main()
